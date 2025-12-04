import os
import json
import re
from datetime import datetime
from typing import Dict, Any, List
import logging

from flask import Flask, render_template, request, jsonify, send_from_directory, abort
from pptx import Presentation
import asyncio
import edge_tts
from moviepy import ImageClip, AudioFileClip, concatenate_videoclips, CompositeVideoClip
import tempfile
from pathlib import Path
import subprocess
import pypdfium2 as pdfium

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
TEMPLATE_DIR = os.path.join(BASE_DIR, "template")
TEMPLATE_PPTX = os.path.join(TEMPLATE_DIR, "template.pptx")
METADATA_JSON = os.path.join(TEMPLATE_DIR, "layout_metadata.json")
OUTPUT_DIR = os.path.join(BASE_DIR, "generated")

os.makedirs(OUTPUT_DIR, exist_ok=True)

app = Flask(__name__)



def load_metadata() -> Dict[str, Any]:
    with open(METADATA_JSON, "r", encoding="utf-8") as f:
        meta = json.load(f)
    # Precompute helpers
    layouts_by_name = {l["layout_name"].strip().lower(): l for l in meta.get("layouts", [])}
    must_have_names = [
        l["layout_name"] for l in meta.get("layouts", [])
        if "must have" in l.get("layout_description", "").lower()
    ]
    ignore_names = [
        l["layout_name"] for l in meta.get("layouts", [])
        if "ignore" in l.get("layout_description", "").lower()
    ]
    meta["_layouts_by_name"] = layouts_by_name
    meta["_must_have_names"] = must_have_names
    meta["_ignore_names"] = ignore_names
    return meta


METADATA = load_metadata()


def sentence_case(s: str) -> str:
    s = s.strip()
    if not s:
        return s
    return s[0].upper() + s[1:]


def apply_content_rule(text: str, rule_desc: str) -> str:
    d = rule_desc.lower()
    t = text or ""
    if "all-caps" in d or "all caps" in d:
        t = t.upper()
    elif "title case" in d:
        t = t.title()
    elif "sentence case" in d:
        # basic sentence case
        t = sentence_case(t)
    return t


def clip_text(text: str, max_chars: int | None) -> str:
    if max_chars and max_chars > 0 and len(text) > max_chars:
        return text[:max_chars]
    return text


def find_layout(prs: Presentation, layout_name: str):
    target = (layout_name or "").strip().lower()
    # exact match
    for l in prs.slide_layouts:
        if (l.name or "").strip().lower() == target:
            return l
    # relaxed contains
    for l in prs.slide_layouts:
        if target and target in (l.name or "").strip().lower():
            return l
    # fallback: first layout
    return prs.slide_layouts[0]


def fill_placeholders(slide, layout_name: str, items: Dict[str, str]):
    layout = METADATA["_layouts_by_name"].get(layout_name.strip().lower())
    rule_by_pid: Dict[int, Dict[str, Any]] = {}
    if layout:
        for ph in layout.get("placeholders", []):
            rule_by_pid[ph["id"]] = ph
    # Fill only matching placeholders by placeholder idx
    # Use slide.placeholders to iterate known placeholder shapes
    for shp in slide.placeholders:
        try:
            pid = shp.placeholder_format.idx
        except Exception:
            continue
        key = str(pid)
        if key in items and hasattr(shp, "text_frame"):
            text = items[key]
            rule_desc = rule_by_pid.get(pid, {}).get("content_description", "")
            # maxchars = rule_by_pid.get(pid, {}).get("maxchars")  # Not used anymore
            text = apply_content_rule(text, rule_desc)
            # text = clip_text(text, maxchars)  # DISABLED: No longer clipping text to allow full content
            tf = shp.text_frame
            tf.clear()
            tf.text = text


def safe_slug(s: str) -> str:
    s = s.strip().lower()
    s = re.sub(r"[^a-z0-9\-\_\s]", "", s)
    s = re.sub(r"\s+", "-", s)
    return s or "presentation"


def call_gemini_for_plan(topic: str, metadata: Dict[str, Any]) -> Dict[str, Any]:
    api_key = os.environ.get("GEMINI_API_KEY")
    if not api_key:
        app.logger.error("GEMINI_API_KEY not found in environment variables.")
        return {}

    try:
        from google import genai
        from google.genai import types
        
        # Initialize the client with API key
        client = genai.Client(api_key=api_key)
        
        system_instruction = (
            "You are an expert presentation creator for educational content. Your task is to generate a comprehensive slide plan in JSON format. "
            "Your output MUST be a JSON object with a 'slides' array. "
            "You will be given a topic and layout metadata with specific usage rules. "
            "CRITICAL USAGE RULES:\\n"
            "1. MUST HAVE layouts: These MUST be included in your plan.\\n"
            "2. 'Use ONLY ONCE' layouts: Can only appear once in the entire presentation, even if optional.\\n"
            "3. 'Can be used MULTIPLE TIMES' layouts: You can use these as many times as needed to fully explain the topic.\\n"
            "4. Mutually exclusive layouts: If a layout says 'Either use this OR layout_id X', choose only one of them.\\n"
            "5. IGNORE layouts: Never use these.\\n"
            "SLIDE COUNT: Aim for 15-27 slides (can vary by 1-2 slides if absolutely needed for proper completion). "
            "Use repeatable layouts multiple times for different sub-topics, concepts, examples, and explanations. "
            "The goal is COMPLETE and THOROUGH educational coverage. "
            "SPEECH GENERATION: For EACH slide, you MUST generate a 'speech' field. "
            "The speech should be a natural, engaging, and explanatory narration script (approx. 30-60 seconds per slide). "
            "DO NOT just read the text on the slide. Instead, explain the concepts, give examples, and use transitional phrases like 'Moving on to...', 'As we can see here...', 'This is classified into...'. "
            "Make it sound like a real teacher explaining the topic. "
            "Your entire response must be ONLY the raw JSON, without any other text or markdown formatting."
        )
        
        schema_hint = {
            "slides": [
                {
                    "layout_name": "Blank",
                    "placeholders": {"10": "TITLE IN ALL CAPS", "11": "Title Case subtitle"},
                    "speech": "Welcome to this presentation on [Topic]. Today we will explore..."
                }
            ]
        }
        
        prompt = (
            f"Generate a comprehensive JSON slide plan for an educational presentation on the topic: '{topic}'.\\n\\n"
            f"CRITICAL INSTRUCTIONS:\\n"
            f"1. Create 15-27 slides (can vary by 1-2 slides if needed for proper completion). This is the target range.\\n"
            f"2. THOROUGHLY cover all important aspects, sub-topics, concepts, examples, and details of '{topic}'.\\n"
            f"3. Use repeatable layouts MULTIPLE times for different content (e.g., use layout_id 2 for every sub-topic, layout_id 4,5,7,8,9 multiple times for different explanations).\\n"
            f"4. The presentation should be COMPREHENSIVE and EDUCATIONAL.\\n"
            f"5. GENERATE SPEECH: Include a 'speech' field for every slide. This is the narration script. It must be natural, explanatory, and NOT just reading the slide text. Use connecting phrases.\\n\\n"
            f"Layout Usage Rules:\\n"
            f"- Use ALL layouts marked 'MUST HAVE' exactly once\\n"
            f"- For layouts marked 'Use ONLY ONCE', include them at most once\\n"
            f"- For layouts marked 'Can be used MULTIPLE TIMES', use them AS MANY TIMES AS NEEDED for different sub-topics/concepts\\n"
            f"- For mutually exclusive layouts (marked 'Either use this OR layout_id X'), choose only one\\n"
            f"- Never use layouts marked 'IGNORE'\\n\\n"
            f"Layout Metadata:\\n{json.dumps(metadata, ensure_ascii=False)}\\n\\n"
            f"Target: 15-27 slides with engaging narration scripts.\\n"
            f"Respond with only the raw JSON output, matching this schema: {json.dumps(schema_hint)}\\n"
        )
        
        app.logger.info(f"Sending prompt to Gemini for topic: '{topic}'")
        app.logger.debug(f"System Instruction: {system_instruction}")
        app.logger.debug(f"User Prompt: {prompt}")

        # Use the new genai API with system instruction
        response = client.models.generate_content(
            model='gemini-2.5-flash',
            contents=prompt,
            config=types.GenerateContentConfig(
                system_instruction=system_instruction,
                temperature=0.7,
            )
        )
        
        text = response.text if hasattr(response, 'text') else None
        
        app.logger.info(f"Raw response from Gemini:\n{text}")

        if not text:
            app.logger.error("Received empty response from Gemini.")
            return {}
            
        # Extract JSON - handle markdown code blocks
        # Remove markdown code block markers if present
        text_clean = text.strip()
        if text_clean.startswith('```'):
            # Remove opening ```json or ```
            text_clean = re.sub(r'^```(?:json)?\s*\n?', '', text_clean)
            # Remove closing ```
            text_clean = re.sub(r'\n?```\s*$', '', text_clean)
        
        # Now extract the JSON object
        m = re.search(r'\{[\s\S]*\}', text_clean.strip())
        raw = m.group(0) if m else text_clean
        
        # Sanitize JSON: Replace curly quotes and other problematic characters
        raw = raw.replace('"', '"').replace('"', '"')  # Curly double quotes to straight
        raw = raw.replace(''', "'").replace(''', "'")  # Curly single quotes to straight
        raw = raw.replace('…', '...')  # Ellipsis
        
        app.logger.info(f"Extracted JSON string for parsing:\n{raw}")
        
        plan = json.loads(raw)
        
        app.logger.info(f"Successfully parsed JSON plan. Number of slides: {len(plan.get('slides', []))}")
        
        return plan if isinstance(plan, dict) else {}
    except Exception as e:
        app.logger.error(f"An error occurred in call_gemini_for_plan: {e}", exc_info=True)
        return {}




def build_pptx_from_plan(topic: str, plan: Dict[str, Any]) -> str:
    app.logger.info(f"Starting PPTX build for topic: '{topic}'")
    prs = Presentation(TEMPLATE_PPTX)

    # Build slides from plan
    slides_to_build = plan.get("slides", [])
    app.logger.info(f"Building {len(slides_to_build)} slides from the plan.")
    for i, slide_spec in enumerate(slides_to_build):
        layout_name = slide_spec.get("layout_name")
        placeholders = slide_spec.get("placeholders", {})
        speech = slide_spec.get("speech", "")
        app.logger.info(f"Building slide {i+1}/{len(slides_to_build)} with layout: '{layout_name}'")
        if not layout_name:
            app.logger.warning(f"Skipping slide {i+1} due to missing layout name.")
            continue
        layout = find_layout(prs, layout_name)
        slide = prs.slides.add_slide(layout)
        
        # Save speech to notes
        if speech:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = speech
            
        try:
            fill_placeholders(slide, layout_name, placeholders)
        except Exception as e:
            # Continue even if a placeholder can't be filled
            app.logger.error(f"Failed to fill placeholder on slide {i+1}: {e}", exc_info=True)
            pass

    slug = safe_slug(topic)
    fname = f"{slug}-{datetime.now().strftime('%Y%m%d-%H%M%S')}.pptx"
    out_path = os.path.join(OUTPUT_DIR, fname)
    prs.save(out_path)
    app.logger.info(f"Successfully saved presentation to {out_path}")
    return fname


@app.route("/")
def index():
    return render_template("index.html")


@app.route("/generate", methods=["POST"])
def generate():
    app.logger.info("Received request to /generate")
    try:
        data = request.get_json(force=True)
        topic = (data.get("topic") or "").strip()
        app.logger.info(f"Request topic: '{topic}'")
    except Exception as e:
        app.logger.error(f"Error parsing request JSON: {e}", exc_info=True)
        return jsonify({"error": "Invalid request"}), 400

    if not topic:
        app.logger.warning("Topic is required but was not provided.")
        return jsonify({"error": "Topic is required"}), 400

    plan = call_gemini_for_plan(topic, METADATA)
    if not plan or not plan.get("slides"):
        app.logger.error("The AI failed to generate a valid presentation plan.")
        return jsonify({"error": "The AI failed to generate a valid presentation plan. Please try again."}), 500

    try:
        filename = build_pptx_from_plan(topic, plan)
    except FileNotFoundError:
        app.logger.error(f"Template PPTX not found at {TEMPLATE_PPTX}")
        return jsonify({"error": "Template PPTX not found."}), 500
    except Exception as e:
        app.logger.error(f"Failed to build PPTX: {e}", exc_info=True)
        return jsonify({"error": f"Failed to build PPTX: {e}"}), 500

    app.logger.info(f"Successfully generated presentation: {filename}")
    return jsonify({
        "filename": filename,
        "download_url": f"/download/{filename}"
    })


@app.route("/download/<path:filename>")
def download(filename):
    if not os.path.exists(os.path.join(OUTPUT_DIR, filename)):
        abort(404)
    return send_from_directory(OUTPUT_DIR, filename, as_attachment=True, mimetype=(
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ))


@app.route('/generate_video', methods=['POST'])
def generate_video():
    data = request.json
    filename = data.get('filename')
    if not filename:
        return jsonify({"error": "Filename is required"}), 400
        
    pptx_path = os.path.join(OUTPUT_DIR, filename)
    if not os.path.exists(pptx_path):
        return jsonify({"error": "File not found"}), 404
        
    video_filename = filename.replace(".pptx", ".mp4")
    video_path = os.path.join(OUTPUT_DIR, video_filename)
    
    try:
        # Run async video generation in a sync wrapper
        asyncio.run(create_video_from_pptx(pptx_path, video_path))
        return jsonify({"video_url": f"/download/{video_filename}"})
    except Exception as e:
        app.logger.error(f"Video generation failed: {e}", exc_info=True)
        return jsonify({"error": str(e)}), 500

async def create_video_from_pptx(pptx_path, output_video_path):
    app.logger.info(f"Starting video generation for {pptx_path}")
    
    with tempfile.TemporaryDirectory() as temp_dir:
        # 1. Convert PPTX to PDF using C# app
        pdf_path = os.path.join(temp_dir, "slides.pdf")
        converter_dll = os.path.join(BASE_DIR, "csharp_converter", "bin", "Debug", "net9.0", "csharp_converter.dll")
        
        # Check if DLL exists
        if not os.path.exists(converter_dll):
            raise FileNotFoundError(f"C# Converter DLL not found at {converter_dll}. Please build the C# project.")
            
        cmd = ["dotnet", converter_dll, pptx_path, pdf_path]
        app.logger.info(f"Running conversion command: {' '.join(cmd)}")
        subprocess.run(cmd, check=True, capture_output=True)
        
        # 2. Convert PDF to Images
        app.logger.info("Converting PDF to images...")
        pdf = pdfium.PdfDocument(pdf_path)
        images = []
        for i in range(len(pdf)):
            page = pdf[i]
            bitmap = page.render(scale=2) # 2x scale for better quality
            pil_image = bitmap.to_pil()
            img_path = os.path.join(temp_dir, f"slide_{i}.png")
            pil_image.save(img_path)
            images.append(img_path)
            
        # 3. Generate Audio and Video Clips
        app.logger.info("Generating audio and video clips...")
        prs = Presentation(pptx_path)
        video_clips = []
        
        # Ensure we don't process more slides than we have images for
        num_slides = min(len(prs.slides), len(images))
        
        for i in range(num_slides):
            slide = prs.slides[i]
            
            # Extract notes
            notes = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text
            
            if not notes or len(notes.strip()) < 5:
                notes = f"Slide {i+1}. Please review the content on this slide."
            
            # Generate Audio
            audio_path = os.path.join(temp_dir, f"audio_{i}.mp3")
            communicate = edge_tts.Communicate(notes, "en-US-AriaNeural")
            await communicate.save(audio_path)
            
            # Create Clips
            audio_clip = AudioFileClip(audio_path)
            img_clip = ImageClip(images[i]).set_duration(audio_clip.duration)
            
            # Add a slight fade in/out for smoother transitions
            img_clip = img_clip.fadein(0.5).fadeout(0.5)
            
            img_clip = img_clip.set_audio(audio_clip)
            video_clips.append(img_clip)
            
        # 4. Concatenate and Save
        app.logger.info("Rendering final video...")
        final_video = concatenate_videoclips(video_clips, method="compose")
        final_video.write_videofile(output_video_path, fps=24, codec="libx264", audio_codec="aac")
        app.logger.info(f"Video saved to {output_video_path}")


if __name__ == "__main__":
    # Set up Flask's logger
    handler = logging.StreamHandler()
    handler.setLevel(logging.INFO)
    formatter = logging.Formatter('%(asctime)s - %(name)s - %(levelname)s - %(message)s')
    handler.setFormatter(formatter)
    app.logger.addHandler(handler)
    app.logger.setLevel(logging.INFO)
    
    app.run(host="0.0.0.0", port=5000, debug=False)
